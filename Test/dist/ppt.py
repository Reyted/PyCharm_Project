from pptx import Presentation

def read_ppt_to_console(file_path):
    try:
        # 打开PPT文件
        prs = Presentation(file_path)
        print(f"Number of slides: {len(prs.slides)}")

        # 遍历每一页幻灯片并读取其内容
        for slide_index, slide in enumerate(prs.slides, start=1):
            print(f"\nSlide {slide_index}:")

            # 遍历幻灯片中的每一个形状
            for shape in slide.shapes:
                # 如果形状有文本，则输出文本内容
                if hasattr(shape, "text"):
                    print(f"  {shape.text}")
                # 还可以检查其他类型的形状，如图片、图表等，并根据需要处理

    except Exception as e:
        print(f"An error occurred: {e}")

# 调用函数并传递PPT文件路径
ppt_file_path = "C:/Users/24253/Desktop/Python/Python_PPT/333.pptx"
read_ppt_to_console(ppt_file_path)
